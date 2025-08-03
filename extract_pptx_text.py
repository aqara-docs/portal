import os
from pptx import Presentation
import re

def extract_text_from_pptx(pptx_path):
    """PowerPoint 파일에서 텍스트를 추출합니다."""
    try:
        prs = Presentation(pptx_path)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_text.append(f"## 슬라이드 {slide_num}")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # 텍스트 정리
                    cleaned_text = shape.text.strip()
                    if cleaned_text:
                        slide_text.append(cleaned_text)
            
            if len(slide_text) > 1:  # 제목 외에 내용이 있는 경우
                text_content.extend(slide_text)
                text_content.append("")  # 빈 줄 추가
        
        return "\n".join(text_content)
    
    except Exception as e:
        return f"텍스트 추출 실패: {str(e)}"

def process_pptx_files():
    """PowerPoint 파일들을 처리합니다."""
    pptx_files = [
        "./pages/rag_files/process/☆ 신규입사자 온보딩 자료_250403.pptx",
        "./pages/rag_files/process/[경영지원_전자서명_006] 전자서명_모두싸인 사용가이드_250331.pptx"
    ]
    
    for pptx_path in pptx_files:
        if os.path.exists(pptx_path):
            print(f"처리 중: {pptx_path}")
            
            # 파일명에서 확장자 제거
            base_name = os.path.splitext(os.path.basename(pptx_path))[0]
            
            # 텍스트 추출
            text_content = extract_text_from_pptx(pptx_path)
            
            # Markdown 파일 생성
            md_path = os.path.join(os.path.dirname(pptx_path), f"{base_name}.md")
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(f"# {base_name}\n\n")
                f.write(f"원본 파일: {os.path.basename(pptx_path)}\n\n")
                f.write("---\n\n")
                f.write(text_content)
            
            print(f"  Markdown 파일 생성: {md_path}")
            
            # 원본 파일을 임시 폴더로 이동
            temp_path = f"./pages/rag_files/temp_large_files/{os.path.basename(pptx_path)}"
            os.rename(pptx_path, temp_path)
            print(f"  원본 파일 이동: {temp_path}")
            
            # 토큰 수 추정
            estimated_tokens = len(text_content) // 4
            print(f"  추정 토큰 수: {estimated_tokens:,}")
            print()
        else:
            print(f"파일 없음: {pptx_path}")

if __name__ == "__main__":
    process_pptx_files() 