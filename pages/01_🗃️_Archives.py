import streamlit as st
import mysql.connector
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import os
from dotenv import load_dotenv
from datetime import datetime, timedelta
import json
import base64
import io
from io import StringIO
from openai import OpenAI
import anthropic
import requests
import traceback
import re
from typing import Dict, List, Any, Optional
import tempfile
import PyPDF2
import docx
from pptx import Presentation
import openpyxl
import pandas as pd
from langchain_anthropic import ChatAnthropic
import time
import math
from PyPDF2 import PdfReader, PdfWriter

# 환경 변수 로드
load_dotenv()

# 페이지 설정
st.set_page_config(
    page_title="📁 ARCHIVES",
    page_icon="📁",
    layout="wide",
    initial_sidebar_state="expanded"
)

st.title("📁 ARCHIVES")
# 인증 기능 (간단한 비밀번호 보호)
if 'authenticated' not in st.session_state:
    st.session_state.authenticated = False

admin_pw = os.getenv('ADMIN_PASSWORD')
if not admin_pw:
    st.error('환경변수(ADMIN_PASSWORD)가 설정되어 있지 않습니다. .env 파일을 확인하세요.')
    st.stop()

if not st.session_state.authenticated:
    password = st.text_input("관리자 비밀번호를 입력하세요", type="password")
    if password == admin_pw:
        st.session_state.authenticated = True
        st.rerun()
    else:
        if password:  # 비밀번호가 입력된 경우에만 오류 메시지 표시
            st.error("관리자 권한이 필요합니다")
        st.stop()

def connect_to_db():
    """데이터베이스 연결"""
    try:
        conn = mysql.connector.connect(
            host=os.getenv('SQL_HOST'),
            user=os.getenv('SQL_USER'),
            password=os.getenv('SQL_PASSWORD'),
            database=os.getenv('SQL_DATABASE_NEWBIZ'),
            charset='utf8mb4',
            collation='utf8mb4_unicode_ci'
        )
        return conn
    except mysql.connector.Error as err:
        st.error(f"데이터베이스 연결 오류: {err}")
        return None

def create_file_storage_tables():
    """파일 저장소 테이블 생성"""
    conn = connect_to_db()
    if not conn:
        return False
    
    cursor = conn.cursor()
    
    try:
        # 파일 저장소 메인 테이블
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS file_storage (
                file_id INT AUTO_INCREMENT PRIMARY KEY,
                title VARCHAR(500) NOT NULL,
                author VARCHAR(200),
                description TEXT,
                file_count INT DEFAULT 0,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP ON UPDATE CURRENT_TIMESTAMP,
                INDEX idx_title (title),
                INDEX idx_author (author),
                INDEX idx_created_at (created_at)
            ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4 COLLATE=utf8mb4_unicode_ci
        """)
        
        # 파일 데이터 테이블
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS file_storage_files (
                file_data_id INT AUTO_INCREMENT PRIMARY KEY,
                storage_id INT,
                filename VARCHAR(500) NOT NULL,
                file_type VARCHAR(50),
                file_content LONGTEXT,
                file_binary_data LONGBLOB,
                file_size BIGINT,
                uploaded_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (storage_id) REFERENCES file_storage(file_id) ON DELETE CASCADE,
                INDEX idx_storage_id (storage_id),
                INDEX idx_file_type (file_type),
                INDEX idx_filename (filename)
            ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4 COLLATE=utf8mb4_unicode_ci
        """)
        
        # AI 분석 결과 테이블
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS file_storage_ai_analysis (
                analysis_id INT AUTO_INCREMENT PRIMARY KEY,
                storage_id INT,
                model_name VARCHAR(100),
                analysis_content LONGTEXT,
                summary TEXT,
                key_points TEXT,
                recommendations TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (storage_id) REFERENCES file_storage(file_id) ON DELETE CASCADE,
                INDEX idx_storage_id (storage_id),
                INDEX idx_model_name (model_name)
            ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4 COLLATE=utf8mb4_unicode_ci
        """)
        
        conn.commit()
        return True
        
    except mysql.connector.Error as err:
        st.error(f"테이블 생성 오류: {err}")
        conn.rollback()
        return False
    finally:
        cursor.close()
        conn.close()

def parse_uploaded_file(uploaded_file):
    """업로드된 파일 파싱"""
    try:
        file_extension = uploaded_file.name.split('.')[-1].lower()
        content = ""
        
        uploaded_file.seek(0)
        binary_data = uploaded_file.read()
        binary_base64 = base64.b64encode(binary_data).decode('utf-8')
        
        uploaded_file.seek(0)
        
        if file_extension == 'pdf':
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                content += page.extract_text() + "\n"
                
        elif file_extension == 'docx':
            doc = docx.Document(uploaded_file)
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"
        
        elif file_extension == 'pptx':
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
        
        elif file_extension == 'xlsx':
            workbook = openpyxl.load_workbook(uploaded_file, data_only=True)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content += f"=== {sheet_name} ===\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip():
                        content += row_text + "\n"
                content += "\n"
                
        elif file_extension in ['txt', 'md']:
            uploaded_file.seek(0)
            content = uploaded_file.read().decode('utf-8')
            
        elif file_extension in ['jpg', 'jpeg', 'png', 'gif']:
            content = f"[{file_extension.upper()} 이미지 파일 - {uploaded_file.name}]"
            
        else:
            try:
                uploaded_file.seek(0)
                content = uploaded_file.read().decode('utf-8', errors='ignore')
                if not content.strip():
                    content = f"[{file_extension.upper()} 파일 - 텍스트 추출 불가]"
            except:
                content = f"[{file_extension.upper()} 파일 - 텍스트 추출 불가]"
            
        return {
            'filename': uploaded_file.name,
            'file_type': file_extension,
            'content': content,
            'binary_data': binary_base64,
            'size': len(binary_data)
        }
        
    except Exception as e:
        st.error(f"파일 파싱 오류: {str(e)}")
        return None

def save_file_storage(storage_data):
    """파일 저장소 메인 정보 저장"""
    conn = connect_to_db()
    if not conn:
        return None
    
    cursor = conn.cursor()
    
    try:
        cursor.execute("""
            INSERT INTO file_storage 
            (title, author, description, file_count)
            VALUES (%s, %s, %s, %s)
        """, (
            storage_data['title'],
            storage_data['author'],
            storage_data['description'],
            storage_data['file_count']
        ))
        
        storage_id = cursor.lastrowid
        conn.commit()
        return storage_id
        
    except mysql.connector.Error as err:
        st.error(f"파일 저장소 저장 오류: {err}")
        conn.rollback()
        return None
    finally:
        cursor.close()
        conn.close()

def save_file_data(storage_id, file_data):
    """파일 데이터 저장"""
    conn = connect_to_db()
    if not conn:
        return False
    
    cursor = conn.cursor()
    
    try:
        cursor.execute("""
            INSERT INTO file_storage_files 
            (storage_id, filename, file_type, file_content, file_binary_data, file_size)
            VALUES (%s, %s, %s, %s, %s, %s)
        """, (
            storage_id,
            file_data['filename'],
            file_data['file_type'],
            file_data['content'],
            file_data['binary_data'],
            file_data['size']
        ))
        
        conn.commit()
        return True
        
    except mysql.connector.Error as err:
        st.error(f"파일 데이터 저장 오류: {err}")
        conn.rollback()
        return False
    finally:
        cursor.close()
        conn.close()

def save_ai_analysis(storage_id, model_name, analysis_content, summary, key_points, recommendations):
    """AI 분석 결과 저장"""
    conn = connect_to_db()
    if not conn:
        return False
    
    cursor = conn.cursor()
    
    try:
        cursor.execute("""
            INSERT INTO file_storage_ai_analysis 
            (storage_id, model_name, analysis_content, summary, key_points, recommendations)
            VALUES (%s, %s, %s, %s, %s, %s)
        """, (storage_id, model_name, analysis_content, summary, key_points, recommendations))
        
        conn.commit()
        return True
        
    except mysql.connector.Error as err:
        st.error(f"AI 분석 저장 오류: {err}")
        conn.rollback()
        return False
    finally:
        cursor.close()
        conn.close()

def get_file_storages(search_term=None):
    """파일 저장소 목록 조회"""
    conn = connect_to_db()
    if not conn:
        return []
    
    cursor = conn.cursor(dictionary=True)
    
    try:
        if search_term:
            cursor.execute("""
                SELECT * FROM file_storage 
                WHERE title LIKE %s OR author LIKE %s OR description LIKE %s
                ORDER BY created_at DESC
            """, (f"%{search_term}%", f"%{search_term}%", f"%{search_term}%"))
        else:
            cursor.execute("""
                SELECT * FROM file_storage 
                ORDER BY created_at DESC
            """)
        return cursor.fetchall()
    finally:
        cursor.close()
        conn.close()

def get_storage_files(storage_id):
    """특정 저장소의 파일 목록 조회"""
    conn = connect_to_db()
    if not conn:
        return []
    
    cursor = conn.cursor(dictionary=True)
    
    try:
        cursor.execute("""
            SELECT * FROM file_storage_files 
            WHERE storage_id = %s
            ORDER BY uploaded_at
        """, (storage_id,))
        return cursor.fetchall()
    finally:
        cursor.close()
        conn.close()

def get_ai_analysis(storage_id):
    """AI 분석 결과 조회"""
    conn = connect_to_db()
    if not conn:
        return []
    
    cursor = conn.cursor(dictionary=True)
    
    try:
        cursor.execute("""
            SELECT * FROM file_storage_ai_analysis 
            WHERE storage_id = %s
            ORDER BY created_at DESC
        """, (storage_id,))
        return cursor.fetchall()
    finally:
        cursor.close()
        conn.close()

def get_file_binary_data(file_data_id):
    """파일의 원본 바이너리 데이터 조회"""
    conn = connect_to_db()
    if not conn:
        return None
    
    cursor = conn.cursor(dictionary=True)
    
    try:
        cursor.execute("""
            SELECT filename, file_type, file_binary_data, file_size
            FROM file_storage_files 
            WHERE file_data_id = %s
        """, (file_data_id,))
        result = cursor.fetchone()
        
        if result and result['file_binary_data']:
            binary_data = base64.b64decode(result['file_binary_data'])
            return {
                'filename': result['filename'],
                'file_type': result['file_type'],
                'binary_data': binary_data,
                'file_size': result['file_size']
            }
        return None
        
    except mysql.connector.Error as err:
        st.error(f"파일 데이터 조회 오류: {err}")
        return None
    finally:
        cursor.close()
        conn.close()

def ai_analyze_files(files_content, model_name):
    """AI를 사용한 파일 분석"""
    content_text = "\n\n".join([f"파일명: {f['filename']}\n내용:\n{f['content']}" for f in files_content])
    
    prompt = f"""
    다음 파일들을 분석하여 요약, 핵심 포인트, 추천사항을 제공해 주세요:

    {content_text}

    다음 형식으로 답변해 주세요:
    
    ## 📋 전체 요약
    [파일들의 전반적인 내용과 목적을 요약]

    ## 🎯 핵심 포인트
    - [핵심 포인트 1]
    - [핵심 포인트 2]
    - [핵심 포인트 3]
    
    ## 💡 추천사항
    - [추천사항 1]
    - [추천사항 2]
    - [추천사항 3]
    """
    
    try:
        if model_name.startswith('claude'):
            client = ChatAnthropic(
                model=model_name, 
                api_key=os.getenv('ANTHROPIC_API_KEY'), 
                temperature=0.3, 
                max_tokens=4000
            )
            response = client.invoke([
                {"role": "user", "content": prompt}
            ])
            analysis = response.content if hasattr(response, 'content') else str(response)
        else:
            client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
            response = client.chat.completions.create(
                model=model_name,
                messages=[
                    {"role": "system", "content": "당신은 문서 분석 전문가입니다. 주어진 파일들을 분석하여 요약, 핵심 포인트, 추천사항을 제공합니다."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=4000,
                temperature=0.3
            )
            analysis = response.choices[0].message.content
        
        # 분석 결과에서 각 섹션 추출
        sections = analysis.split('##')
        summary = ""
        key_points = ""
        recommendations = ""
        
        for section in sections:
            if '요약' in section:
                summary = section.replace('요약', '').strip()
            elif '핵심 포인트' in section or '핵심' in section:
                key_points = section.replace('핵심 포인트', '').replace('핵심', '').strip()
            elif '추천' in section or '권장' in section:
                recommendations = section.replace('추천사항', '').replace('추천', '').replace('권장', '').strip()
        
        return {
            'full_analysis': analysis,
            'summary': summary,
            'key_points': key_points,
            'recommendations': recommendations
        }
        
    except Exception as e:
        st.error(f"AI 분석 중 오류 발생: {str(e)}")
        return None

def display_file_preview(file_data, file_type, filename):
    """파일 미리보기 표시"""
    try:
        file_type_lower = file_type.lower()
        
        if file_type_lower in ['jpg', 'jpeg', 'png', 'gif']:
            if file_data.get('binary_data'):
                st.image(
                    file_data['binary_data'],
                    caption=f"🖼️ {filename}",
                    use_container_width=True
                )
                return True
        
        elif file_type_lower == 'pdf':
            if file_data.get('binary_data'):
                import math
                from PyPDF2 import PdfReader, PdfWriter
                import tempfile, os

                total_size = file_data['file_size']
                # 1MB 초과시 페이징 미리보기 (1페이지씩)
                if total_size > 1 * 1024 * 1024:
                    key = f"pdf_preview_page_start_{filename}"
                    if key not in st.session_state:
                        st.session_state[key] = 0  # 0-based index

                    # PDF 전체 페이지 수 구하기
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_in:
                        tmp_in.write(file_data['binary_data'])
                        tmp_in_path = tmp_in.name
                    reader = PdfReader(tmp_in_path)
                    total_pages = len(reader.pages)
                    os.unlink(tmp_in_path)

                    page_start = st.session_state[key]
                    page_end = min(page_start + 1, total_pages)

                    col_prev, col_next = st.columns([1, 1])
                    with col_prev:
                        if st.button("⬅️ 이전", disabled=page_start == 0, key=f"prev_{filename}"):
                            st.session_state[key] = max(0, page_start - 1)
                            st.rerun()
                    with col_next:
                        if st.button("다음 ➡️", disabled=page_end >= total_pages, key=f"next_{filename}"):
                            st.session_state[key] = min(total_pages - 1, page_start + 1)
                            st.rerun()

                    # 미리보기 PDF 생성 (1페이지)
                    try:
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_in:
                            tmp_in.write(file_data['binary_data'])
                            tmp_in_path = tmp_in.name
                        reader = PdfReader(tmp_in_path)
                        writer = PdfWriter()
                        for i in range(page_start, page_end):
                            writer.add_page(reader.pages[i])
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                            writer.write(tmp_out)
                            tmp_out_path = tmp_out.name
                        with open(tmp_out_path, "rb") as f:
                            preview_pdf_bytes = f.read()
                        os.unlink(tmp_in_path)
                        os.unlink(tmp_out_path)
                        st.markdown(f"**페이지 {page_start+1} / {total_pages}**")
                        pdf_base64 = base64.b64encode(preview_pdf_bytes).decode('utf-8')
                        pdf_display = f"""
                        <iframe src=\"data:application/pdf;base64,{pdf_base64}\" 
                                width=\"100%\" height=\"600px\" type=\"application/pdf\">
                            <p>PDF를 표시할 수 없습니다. 
                            <a href=\"data:application/pdf;base64,{pdf_base64}\" target=\"_blank\">
                            여기를 클릭하여 새 탭에서 열어보세요.</a></p>
                        </iframe>
                        """
                        st.markdown(pdf_display, unsafe_allow_html=True)
                    except Exception as e:
                        st.error(f"PDF 미리보기 생성 중 오류: {e}")

                    st.download_button(
                        label="💾 전체 PDF 다운로드",
                        data=file_data['binary_data'],
                        file_name=filename,
                        mime="application/pdf"
                    )
                    return True
                else:
                    st.subheader(f"📄 PDF 미리보기: {filename}")
                    pdf_base64 = base64.b64encode(file_data['binary_data']).decode('utf-8')
                    pdf_display = f"""
                    <iframe src=\"data:application/pdf;base64,{pdf_base64}\" 
                            width=\"100%\" height=\"600px\" type=\"application/pdf\">
                        <p>PDF를 표시할 수 없습니다. 
                        <a href=\"data:application/pdf;base64,{pdf_base64}\" target=\"_blank\">
                        여기를 클릭하여 새 탭에서 열어보세요.</a></p>
                    </iframe>
                    """
                    st.markdown(pdf_display, unsafe_allow_html=True)
                    return True
        
        elif file_type_lower in ['txt', 'md']:
            if file_data.get('binary_data'):
                try:
                    text_content = file_data['binary_data'].decode('utf-8')
                    st.subheader(f"📄 텍스트 파일 미리보기: {filename}")
                    
                    if file_type_lower == 'md':
                        st.markdown(text_content)
                    else:
                        st.text_area(
                            "파일 내용",
                            value=text_content,
                            height=400,
                            disabled=True
                        )
                    return True
                except Exception as e:
                    st.error(f"텍스트 파일을 읽는 중 오류 발생: {str(e)}")
                    return False
        
        elif file_type_lower == 'xlsx':
            return display_excel_preview(file_data, filename)
        
        elif file_type_lower == 'pptx':
            return display_powerpoint_preview(file_data, filename)
        
        elif file_type_lower == 'docx':
            return display_word_preview(file_data, filename)
        
        return False
        
    except Exception as e:
        st.error(f"파일 미리보기 오류: {str(e)}")
        return False

def display_excel_preview(file_data, filename):
    """엑셀 파일 미리보기"""
    try:
        st.subheader(f"📊 Excel 미리보기: {filename}")
        
        if file_data.get('binary_data'):
            # 바이너리 데이터를 임시 파일로 저장
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
                tmp_file.write(file_data['binary_data'])
                tmp_file_path = tmp_file.name
            
            try:
                # 엑셀 파일 읽기
                workbook = openpyxl.load_workbook(tmp_file_path, data_only=True)
                
                # 시트 탭 생성
                sheet_names = workbook.sheetnames
                if len(sheet_names) > 1:
                    selected_sheet = st.selectbox("시트 선택", sheet_names, key=f"sheet_select_{filename}")
                else:
                    selected_sheet = sheet_names[0]
                
                # 선택된 시트 표시
                sheet = workbook[selected_sheet]
                
                # 데이터를 pandas DataFrame으로 변환
                data = []
                for row in sheet.iter_rows(values_only=True):
                    data.append(row)
                
                if data:
                    # 첫 번째 행을 헤더로 사용
                    df = pd.DataFrame(data[1:], columns=data[0])
                    
                    # NaN 값을 빈 문자열로 변경
                    df = df.fillna('')
                    
                    st.write(f"**시트: {selected_sheet}**")
                    st.write(f"행 수: {len(df)}, 열 수: {len(df.columns)}")
                    
                    # 테이블 표시 (최대 100행)
                    if len(df) > 100:
                        st.warning("⚠️ 데이터가 많아서 상위 100행만 표시합니다.")
                        st.dataframe(df.head(100), use_container_width=True)
                    else:
                        st.dataframe(df, use_container_width=True)
                else:
                    st.info("빈 시트입니다.")
                
                workbook.close()
                return True
                
            finally:
                # 임시 파일 삭제
                import os
                if os.path.exists(tmp_file_path):
                    os.unlink(tmp_file_path)
        
        return False
    except Exception as e:
        st.error(f"Excel 미리보기 오류: {str(e)}")
        return False

def display_powerpoint_preview(file_data, filename):
    """PowerPoint 파일 미리보기"""
    try:
        st.subheader(f"📋 PowerPoint 미리보기: {filename}")
        
        if file_data.get('binary_data'):
            # 바이너리 데이터를 임시 파일로 저장
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                tmp_file.write(file_data['binary_data'])
                tmp_file_path = tmp_file.name
            
            try:
                # PowerPoint 파일 읽기
                prs = Presentation(tmp_file_path)
                
                st.write(f"총 슬라이드 수: {len(prs.slides)}")
                st.markdown("---")
                
                # 모든 슬라이드를 순서대로 표시
                for i, slide in enumerate(prs.slides, 1):
                    # 슬라이드 헤더
                    st.markdown(f"## 📋 슬라이드 {i}")
                    
                    slide_content = []
                    slide_titles = []
                    
                    # 슬라이드 내용 추출 (텍스트와 이미지를 위치 순서대로)
                    slide_elements = []  # 텍스트와 이미지를 순서대로 저장
                    
                    # shape의 위치 정보를 기준으로 정렬하여 처리
                    sorted_shapes = []
                    for shape in slide.shapes:
                        try:
                            # shape의 top 위치를 기준으로 정렬
                            top_position = shape.top if hasattr(shape, 'top') else 0
                            sorted_shapes.append((top_position, shape))
                        except:
                            sorted_shapes.append((0, shape))
                    
                    # 위치 순서대로 정렬
                    sorted_shapes.sort(key=lambda x: x[0])
                    
                    for top_pos, shape in sorted_shapes:
                        # 텍스트 처리
                        if hasattr(shape, "text") and shape.text.strip():
                            is_title = False
                            try:
                                # placeholder인지 확인
                                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                                        if shape.placeholder_format.type == 1:  # 제목 placeholder
                                            is_title = True
                                # 또는 간단한 휴리스틱으로 제목 감지
                                elif len(shape.text.strip()) < 100 and '\n' not in shape.text.strip():
                                    is_title = True
                            except:
                                pass
                            
                            if is_title:
                                slide_elements.append({
                                    'type': 'title',
                                    'content': shape.text,
                                    'position': top_pos
                                })
                            else:
                                slide_elements.append({
                                    'type': 'text',
                                    'content': shape.text,
                                    'position': top_pos
                                })
                        
                        # 이미지 처리
                        try:
                            from pptx.enum.shapes import MSO_SHAPE_TYPE
                            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                image = shape.image
                                image_bytes = image.blob
                                slide_elements.append({
                                    'type': 'image',
                                    'content': {
                                        'data': image_bytes,
                                        'filename': getattr(image, 'filename', f'slide_{i}_image_{len([e for e in slide_elements if e["type"] == "image"])+1}')
                                    },
                                    'position': top_pos
                                })
                        except Exception as e:
                            pass
                    
                    # 슬라이드 요소들을 위치 순서대로 표시
                    if slide_elements:
                        for element in slide_elements:
                            if element['type'] == 'title':
                                st.markdown(f"### 🎯 {element['content']}")
                                
                            elif element['type'] == 'text':
                                content = element['content']
                                if content.strip():
                                    # 줄바꿈을 기준으로 나누어서 표시
                                    lines = content.split('\n')
                                    for line in lines:
                                        if line.strip():
                                            st.markdown(f"• {line.strip()}")
                                    st.markdown("")  # 텍스트 블록 사이 여백
                                    
                            elif element['type'] == 'image':
                                img = element['content']
                                st.image(
                                    img['data'], 
                                    caption=f"📷 {img['filename']}", 
                                    use_container_width=True
                                )
                                st.markdown("")  # 이미지 사이 여백
                    else:
                        st.info("이 슬라이드에는 내용이 없습니다.")
                    
                    # 슬라이드 구분선 (마지막 슬라이드가 아닌 경우)
                    if i < len(prs.slides):
                        st.markdown("---")
                        st.markdown("")  # 여백 추가
                
                return True
                
            finally:
                # 임시 파일 삭제
                import os
                if os.path.exists(tmp_file_path):
                    os.unlink(tmp_file_path)
        
        return False
    except Exception as e:
        st.error(f"PowerPoint 미리보기 오류: {str(e)}")
        return False

def display_word_preview(file_data, filename):
    """Word 문서 미리보기"""
    try:
        st.subheader(f"📝 Word 문서 미리보기: {filename}")
        
        if file_data.get('binary_data'):
            # 바이너리 데이터를 임시 파일로 저장
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                tmp_file.write(file_data['binary_data'])
                tmp_file_path = tmp_file.name
            
            try:
                # Word 문서 읽기
                doc = docx.Document(tmp_file_path)
                
                st.write(f"총 단락 수: {len(doc.paragraphs)}")
                
                # 문서 내용 표시
                for i, paragraph in enumerate(doc.paragraphs):
                    if paragraph.text.strip():
                        # 스타일에 따라 다른 형식으로 표시
                        text = paragraph.text.strip()
                        
                        # 제목 스타일 감지 (간단한 휴리스틱)
                        if len(text) < 100 and (text.isupper() or any(keyword in text.lower() for keyword in ['제목', '장', '절', 'chapter', 'section'])):
                            st.markdown(f"### {text}")
                        elif paragraph.style.name.startswith('Heading') if hasattr(paragraph, 'style') else False:
                            st.markdown(f"### {text}")
                        else:
                            st.markdown(text)
                        
                        st.markdown("---")  # 단락 구분선
                
                # 표 내용 표시
                if doc.tables:
                    st.subheader("📊 문서 내 표")
                    for i, table in enumerate(doc.tables, 1):
                        st.write(f"**표 {i}**")
                        
                        # 표 데이터를 DataFrame으로 변환
                        table_data = []
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        
                        if table_data:
                            df = pd.DataFrame(table_data)
                            st.dataframe(df, use_container_width=True)
                
                return True
                
            finally:
                # 임시 파일 삭제
                import os
                if os.path.exists(tmp_file_path):
                    os.unlink(tmp_file_path)
        
        return False
    except Exception as e:
        st.error(f"Word 문서 미리보기 오류: {str(e)}")
        return False

def get_file_mime_type(file_type):
    """파일 타입에 따른 MIME 타입 반환"""
    mime_types = {
        'pdf': 'application/pdf',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'txt': 'text/plain',
        'md': 'text/markdown',
        'jpg': 'image/jpeg',
        'jpeg': 'image/jpeg',
        'png': 'image/png',
        'gif': 'image/gif'
    }
    return mime_types.get(file_type.lower(), 'application/octet-stream')

def get_file_storage_by_id(file_id):
    """특정 파일 저장소 상세 조회"""
    conn = connect_to_db()
    if not conn:
        return None
    
    cursor = conn.cursor(dictionary=True)
    
    try:
        cursor.execute("""
            SELECT * FROM file_storage 
            WHERE file_id = %s
        """, (file_id,))
        return cursor.fetchone()
    finally:
        cursor.close()
        conn.close()

def update_file_storage(file_id, storage_data):
    """파일 저장소 정보 수정"""
    conn = connect_to_db()
    if not conn:
        return False
    
    cursor = conn.cursor()
    
    try:
        cursor.execute("""
            UPDATE file_storage SET
            title = %s, author = %s, description = %s,
            updated_at = CURRENT_TIMESTAMP
            WHERE file_id = %s
        """, (
            storage_data['title'],
            storage_data['author'],
            storage_data['description'],
            file_id
        ))
        
        conn.commit()
        return True
        
    except mysql.connector.Error as err:
        st.error(f"파일 저장소 수정 오류: {err}")
        conn.rollback()
        return False
    finally:
        cursor.close()
        conn.close()

def delete_file_storage(file_id):
    """파일 저장소 삭제"""
    conn = connect_to_db()
    if not conn:
        return False
    
    cursor = conn.cursor()
    
    try:
        # 외래 키 제약으로 인해 관련 파일들도 자동 삭제됨
        cursor.execute("""
            DELETE FROM file_storage 
            WHERE file_id = %s
        """, (file_id,))
        
        conn.commit()
        return True
        
    except mysql.connector.Error as err:
        st.error(f"파일 저장소 삭제 오류: {err}")
        conn.rollback()
        return False
    finally:
        cursor.close()
        conn.close()

def delete_storage_file(file_data_id):
    """개별 파일 삭제"""
    conn = connect_to_db()
    if not conn:
        return False
    
    cursor = conn.cursor()
    
    try:
        # 먼저 storage_id 조회
        cursor.execute("""
            SELECT storage_id FROM file_storage_files 
            WHERE file_data_id = %s
        """, (file_data_id,))
        result = cursor.fetchone()
        
        if not result:
            return False
        
        storage_id = result[0]
        
        # 파일 삭제
        cursor.execute("""
            DELETE FROM file_storage_files 
            WHERE file_data_id = %s
        """, (file_data_id,))
        
        # 파일 개수 업데이트
        cursor.execute("""
            UPDATE file_storage SET 
            file_count = (SELECT COUNT(*) FROM file_storage_files WHERE storage_id = %s)
            WHERE file_id = %s
        """, (storage_id, storage_id))
        
        conn.commit()
        return True
        
    except mysql.connector.Error as err:
        st.error(f"파일 삭제 오류: {err}")
        conn.rollback()
        return False
    finally:
        cursor.close()
        conn.close()

def delete_storage_ai_analysis(storage_id):
    """파일 저장소의 모든 AI 분석 결과 삭제"""
    conn = connect_to_db()
    if not conn:
        return False
    
    cursor = conn.cursor()
    
    try:
        cursor.execute("""
            DELETE FROM file_storage_ai_analysis 
            WHERE storage_id = %s
        """, (storage_id,))
        
        deleted_count = cursor.rowcount
        conn.commit()
        return deleted_count
        
    except mysql.connector.Error as err:
        st.error(f"AI 분석 결과 삭제 오류: {err}")
        conn.rollback()
        return False
    finally:
        cursor.close()
        conn.close()

def search_file_storages(search_term=None, author_filter=None, date_from=None, date_to=None):
    """파일 저장소 고급 검색"""
    conn = connect_to_db()
    if not conn:
        return []
    
    cursor = conn.cursor(dictionary=True)
    
    try:
        query = """
            SELECT fs.*, 
                   COUNT(fsf.file_data_id) as actual_file_count,
                   COUNT(fsa.analysis_id) as analysis_count
            FROM file_storage fs
            LEFT JOIN file_storage_files fsf ON fs.file_id = fsf.storage_id
            LEFT JOIN file_storage_ai_analysis fsa ON fs.file_id = fsa.storage_id
            WHERE 1=1
        """
        params = []
        
        if search_term:
            query += " AND (fs.title LIKE %s OR fs.author LIKE %s OR fs.description LIKE %s)"
            search_pattern = f"%{search_term}%"
            params.extend([search_pattern, search_pattern, search_pattern])
        
        if author_filter and author_filter != "전체":
            query += " AND fs.author = %s"
            params.append(author_filter)
        
        if date_from:
            query += " AND DATE(fs.created_at) >= %s"
            params.append(date_from)
        
        if date_to:
            query += " AND DATE(fs.created_at) <= %s"
            params.append(date_to)
        
        query += " GROUP BY fs.file_id ORDER BY fs.created_at DESC"
        
        cursor.execute(query, params)
        return cursor.fetchall()
        
    except mysql.connector.Error as err:
        st.error(f"파일 저장소 검색 오류: {err}")
        return []
    finally:
        cursor.close()
        conn.close()

def get_all_authors():
    """모든 작성자 목록 조회"""
    conn = connect_to_db()
    if not conn:
        return []
    
    cursor = conn.cursor()
    
    try:
        cursor.execute("""
            SELECT DISTINCT author FROM file_storage 
            WHERE author IS NOT NULL AND author != ''
            ORDER BY author
        """)
        return [row[0] for row in cursor.fetchall()]
        
    except mysql.connector.Error as err:
        st.error(f"작성자 목록 조회 오류: {err}")
        return []
    finally:
        cursor.close()
        conn.close()

def replace_storage_file(file_data_id, new_file_data):
    """기존 파일을 새 파일로 교체"""
    conn = connect_to_db()
    if not conn:
        return False
    
    cursor = conn.cursor()
    
    try:
        cursor.execute("""
            UPDATE file_storage_files SET
            filename = %s, file_type = %s, file_content = %s, 
            file_binary_data = %s, file_size = %s, uploaded_at = CURRENT_TIMESTAMP
            WHERE file_data_id = %s
        """, (
            new_file_data['filename'],
            new_file_data['file_type'],
            new_file_data['content'],
            new_file_data['binary_data'],
            new_file_data['size'],
            file_data_id
        ))
        
        conn.commit()
        return True
        
    except mysql.connector.Error as err:
        st.error(f"파일 교체 오류: {err}")
        conn.rollback()
        return False
    finally:
        cursor.close()
        conn.close()

def add_files_to_storage(storage_id, new_files_data):
    """기존 저장소에 새 파일들 추가"""
    conn = connect_to_db()
    if not conn:
        return False
    
    cursor = conn.cursor()
    
    try:
        # 새 파일들 저장
        for file_data in new_files_data:
            cursor.execute("""
                INSERT INTO file_storage_files 
                (storage_id, filename, file_type, file_content, file_binary_data, file_size)
                VALUES (%s, %s, %s, %s, %s, %s)
            """, (
                storage_id,
                file_data['filename'],
                file_data['file_type'],
                file_data['content'],
                file_data['binary_data'],
                file_data['size']
            ))
        
        # 파일 개수 업데이트
        cursor.execute("""
            UPDATE file_storage SET 
            file_count = (SELECT COUNT(*) FROM file_storage_files WHERE storage_id = %s),
            updated_at = CURRENT_TIMESTAMP
            WHERE file_id = %s
        """, (storage_id, storage_id))
        
        conn.commit()
        return True
        
    except mysql.connector.Error as err:
        st.error(f"파일 추가 오류: {err}")
        conn.rollback()
        return False
    finally:
        cursor.close()
        conn.close()

def main():
    # 테이블 생성 확인
    create_file_storage_tables()
    
    # 사이드바 메뉴
    st.sidebar.title("🗂️ 메뉴")
    
    # 편집 모드인 경우 메뉴를 편집으로 자동 설정
    if st.session_state.get('edit_storage_id'):
        default_menu = "✏️ 저장소 편집"
    else:
        default_menu = "📤 파일 업로드"
    
    menu_options = ["📤 파일 업로드", "📋 파일 목록", "🔍 고급 검색", "✏️ 저장소 편집", "📊 통계"]
    
    try:
        default_index = menu_options.index(default_menu)
    except:
        default_index = 0
    
    menu = st.sidebar.selectbox(
        "기능 선택",
        menu_options,
        index=default_index,
        key="main_menu"
    )
    
    # 메뉴 변경 시 편집 상태 정리 (편집 메뉴가 아닌 경우)
    if menu != "✏️ 저장소 편집" and st.session_state.get('edit_storage_id'):
        del st.session_state.edit_storage_id
        if 'return_to_search' in st.session_state:
            del st.session_state.return_to_search
    
    if menu == "📤 파일 업로드":
        st.header("📤 파일 업로드")
        
        with st.form("file_upload_form"):
            # 기본 정보 입력
            col1, col2 = st.columns(2)
            with col1:
                title = st.text_input("📝 제목", placeholder="파일 저장소의 제목을 입력하세요")
                author = st.text_input("✍️ 작성자", placeholder="작성자명을 입력하세요")
            
            with col2:
                description = st.text_area("📋 개요", placeholder="파일들에 대한 간단한 설명을 입력하세요", height=100)
            
            # 파일 업로드
            uploaded_files = st.file_uploader(
                "📁 파일 선택 (PDF, DOCX, PPTX, XLSX, Markdown, txt, jpeg, png)",
                type=['pdf', 'docx', 'pptx', 'xlsx', 'md', 'txt', 'jpg', 'jpeg', 'png'],
                accept_multiple_files=True,
                help="최대 10개 파일까지 업로드 가능합니다."
            )
            
            # AI 분석 옵션
            col3, col4 = st.columns(2)
            with col3:
                enable_ai = st.checkbox("🤖 AI 분석 활성화", value=True)
            with col4:
                if enable_ai:
                    model_options = [
                        "claude-3-7-sonnet-latest",
                        "claude-3-5-sonnet-latest",
                        "claude-3-5-sonnet-20241022",
                        "claude-3-5-haiku-20241022",
                        "claude-3-opus-20240229",
                        "claude-3-sonnet-20240229", 
                        "claude-3-haiku-20240307",
                        "gpt-4o",
                        "gpt-4o-mini",
                        "gpt-4-turbo",
                        "gpt-3.5-turbo",
                        "o1-preview",
                        "o1-mini"
                    ]
                    
                    model_choice = st.selectbox(
                        "AI 모델 선택",
                        model_options,
                        index=0,  # claude-3-5-sonnet-20241022를 기본값으로
                        help="Claude-3-7-sonnet-latest, Claude-3-5-sonnet-latest와 o1 모델들은 Extended Thinking(Reasoning)을 지원합니다."
                    )
            
            submitted = st.form_submit_button("💾 저장", type="primary")
            
            if submitted:
                if not title:
                    st.error("제목을 입력해 주세요.")
                elif not uploaded_files:
                    st.error("최소 1개 이상의 파일을 업로드해 주세요.")
                else:
                    with st.spinner("파일을 처리 중입니다..."):
                        # 파일 저장소 메인 정보 저장
                        storage_data = {
                            'title': title,
                            'author': author or "Unknown",
                            'description': description,
                            'file_count': len(uploaded_files)
                        }
                        
                        storage_id = save_file_storage(storage_data)
                        
                        if storage_id:
                            # 각 파일 처리 및 저장
                            files_content = []
                            success_count = 0
                            
                            for uploaded_file in uploaded_files:
                                file_data = parse_uploaded_file(uploaded_file)
                                if file_data:
                                    if save_file_data(storage_id, file_data):
                                        files_content.append(file_data)
                                        success_count += 1
                            
                            st.success(f"✅ {success_count}/{len(uploaded_files)} 개 파일이 성공적으로 저장되었습니다!")
                            
                            # AI 분석 수행
                            if enable_ai and files_content:
                                with st.spinner("AI 분석을 수행 중입니다..."):
                                    analysis_result = ai_analyze_files(files_content, model_choice)
                                    
                                    if analysis_result:
                                        save_ai_analysis(
                                            storage_id,
                                            model_choice,
                                            analysis_result['full_analysis'],
                                            analysis_result['summary'],
                                            analysis_result['key_points'],
                                            analysis_result['recommendations']
                                        )
                                        st.success("🤖 AI 분석이 완료되었습니다!")
                                        
                                        # 분석 결과 미리보기
                                        with st.expander("📊 AI 분석 결과 미리보기"):
                                            st.markdown(analysis_result['full_analysis'])
                        else:
                            st.error("파일 저장소 생성에 실패했습니다.")
    
    elif menu == "📋 파일 목록":
        st.header("📋 파일 목록")
        
        storages = get_file_storages()
        
        if not storages:
            st.info("저장된 파일이 없습니다.")
            return
        
        # 전체 화면 파일 미리보기 (expander 밖에서 표시)
        if st.session_state.get('preview_file_id'):
            st.markdown("---")
            st.markdown("## 🔍 파일 미리보기 (전체 화면)")
            
            col_preview_1, col_preview_2 = st.columns([10, 1])
            with col_preview_2:
                if st.button("❌ 닫기", key="close_preview"):
                    del st.session_state.preview_file_id
                    del st.session_state.preview_filename
                    del st.session_state.preview_file_type
                    st.rerun()
            
            with col_preview_1:
                file_data = get_file_binary_data(st.session_state.preview_file_id)
                if file_data:
                    display_file_preview(
                        file_data, 
                        st.session_state.preview_file_type, 
                        st.session_state.preview_filename
                    )
                else:
                    st.error("파일 데이터를 불러올 수 없습니다.")
            
            st.markdown("---")
        
        for storage in storages:
            with st.expander(f"📁 {storage['title']} ({storage['author']}) - {storage['created_at'].strftime('%Y-%m-%d %H:%M')}"):
                col1, col2 = st.columns([3, 1])
                
                with col1:
                    st.write(f"**작성자:** {storage['author']}")
                    st.write(f"**파일 개수:** {storage['file_count']}")
                    if storage['description']:
                        st.write(f"**개요:** {storage['description']}")
                
                with col2:
                    col2_1, col2_2, col2_3 = st.columns(3)
                    with col2_1:
                        if st.button("📂", key=f"view_{storage['file_id']}", help="파일 보기"):
                            st.session_state.selected_storage = storage['file_id']
                    with col2_2:
                        if st.button("✏️", key=f"edit_{storage['file_id']}", help="편집"):
                            st.session_state.edit_storage_id = storage['file_id']
                            st.rerun()
                    with col2_3:
                        if st.button("🗑️", key=f"delete_{storage['file_id']}", help="삭제"):
                            st.session_state.delete_storage_id = storage['file_id']
                
                # 선택된 저장소의 파일들 표시
                if st.session_state.get('selected_storage') == storage['file_id']:
                    files = get_storage_files(storage['file_id'])
                    
                    if files:
                        st.subheader("📄 파일 목록")
                        for file_info in files:
                            col1, col2, col3 = st.columns([2, 1, 1])
                            
                            with col1:
                                st.write(f"**{file_info['filename']}** ({file_info['file_type'].upper()})")
                                st.caption(f"크기: {file_info['file_size']:,} bytes")
                            
                            with col2:
                                if st.button(f"👁️ 미리보기", key=f"preview_{file_info['file_data_id']}"):
                                    st.session_state.preview_file_id = file_info['file_data_id']
                                    st.session_state.preview_filename = file_info['filename']
                                    st.session_state.preview_file_type = file_info['file_type']
                            
                            with col3:
                                # 버튼들을 가로로 배치
                                file_data = get_file_binary_data(file_info['file_data_id'])
                                if file_data:
                                    mime_type = get_file_mime_type(file_data['file_type'])
                                    st.download_button(
                                        label="💾 다운로드",
                                        data=file_data['binary_data'],
                                        file_name=file_data['filename'],
                                        mime=mime_type,
                                        key=f"download_{file_info['file_data_id']}",
                                        use_container_width=True
                                    )
                                
                                if st.button("🔄 교체", key=f"replace_btn_{file_info['file_data_id']}", use_container_width=True):
                                    st.session_state[f'replace_file_{file_info["file_data_id"]}'] = True
                                
                                if st.button("🗑️ 삭제", key=f"delete_file_{file_info['file_data_id']}", use_container_width=True):
                                    if delete_storage_file(file_info['file_data_id']):
                                        st.success("파일이 삭제되었습니다!")
                                        st.rerun()
                                    else:
                                        st.error("파일 삭제에 실패했습니다.")
                            
                            # 파일 교체 UI
                            if st.session_state.get(f'replace_file_{file_info["file_data_id"]}'):
                                st.write("**파일 교체**")
                                replacement_file = st.file_uploader(
                                    f"새 파일 선택 (현재: {file_info['filename']})",
                                    type=['pdf', 'docx', 'pptx', 'xlsx', 'md', 'txt', 'jpg', 'jpeg', 'png'],
                                    key=f"replace_upload_{file_info['file_data_id']}"
                                )
                                
                                col_rep_1, col_rep_2 = st.columns(2)
                                with col_rep_1:
                                    if st.button("✅ 교체 확인", key=f"confirm_replace_{file_info['file_data_id']}"):
                                        if replacement_file:
                                            new_file_data = parse_uploaded_file(replacement_file)
                                            if new_file_data and replace_storage_file(file_info['file_data_id'], new_file_data):
                                                st.success("파일이 교체되었습니다!")
                                                del st.session_state[f'replace_file_{file_info["file_data_id"]}']
                                                st.rerun()
                                            else:
                                                st.error("파일 교체에 실패했습니다.")
                                        else:
                                            st.error("교체할 파일을 선택해 주세요.")
                                with col_rep_2:
                                    if st.button("❌ 취소", key=f"cancel_replace_{file_info['file_data_id']}"):
                                        del st.session_state[f'replace_file_{file_info["file_data_id"]}']
                                        st.rerun()
                    
                    # 새 파일 추가 기능
                    st.subheader("➕ 새 파일 추가")
                    new_files = st.file_uploader(
                        "새 파일 선택",
                        type=['pdf', 'docx', 'pptx', 'xlsx', 'md', 'txt', 'jpg', 'jpeg', 'png'],
                        accept_multiple_files=True,
                        key=f"add_files_{storage['file_id']}"
                    )
                    
                    if new_files and st.button("➕ 파일 추가", key=f"add_files_btn_{storage['file_id']}"):
                        with st.spinner("파일을 추가하는 중..."):
                            new_files_data = []
                            for new_file in new_files:
                                file_data = parse_uploaded_file(new_file)
                                if file_data:
                                    new_files_data.append(file_data)
                            
                            if new_files_data and add_files_to_storage(storage['file_id'], new_files_data):
                                st.success(f"{len(new_files_data)}개 파일이 추가되었습니다!")
                                st.rerun()
                            else:
                                st.error("파일 추가에 실패했습니다.")
                    
                    # AI 분석 결과 표시
                    analyses = get_ai_analysis(storage['file_id'])
                    if analyses:
                        col_ai_1, col_ai_2 = st.columns([4, 1])
                        with col_ai_1:
                            st.subheader("🤖 AI 분석 결과")
                        with col_ai_2:
                            if st.button("🗑️ 분석 삭제", key=f"delete_analysis_{storage['file_id']}", help="모든 AI 분석 결과 삭제"):
                                deleted_count = delete_storage_ai_analysis(storage['file_id'])
                                if deleted_count:
                                    st.success(f"{deleted_count}개의 AI 분석 결과가 삭제되었습니다!")
                                    st.rerun()
                                else:
                                    st.error("AI 분석 결과 삭제에 실패했습니다.")
                        
                        for analysis in analyses:
                            st.write(f"**모델:** {analysis['model_name']}")
                            st.write(f"**분석 시간:** {analysis['created_at'].strftime('%Y-%m-%d %H:%M')}")
                            st.markdown(analysis['analysis_content'])
                            st.divider()
                
                # 삭제 확인 대화상자
                if st.session_state.get('delete_storage_id') == storage['file_id']:
                    st.error("⚠️ 정말로 이 파일 저장소를 삭제하시겠습니까?")
                    st.write("**삭제될 내용:**")
                    st.write(f"- 저장소: {storage['title']}")
                    st.write(f"- 파일: {storage['file_count']}개")
                    
                    # 분석 결과 개수 확인
                    analyses_count = len(get_ai_analysis(storage['file_id']))
                    if analyses_count > 0:
                        st.write(f"- AI 분석 결과: {analyses_count}개")
                    
                    col_del_1, col_del_2, col_del_3 = st.columns(3)
                    with col_del_1:
                        if st.button("✅ 삭제 확인", key=f"confirm_delete_{storage['file_id']}", type="primary"):
                            if delete_file_storage(storage['file_id']):
                                st.success("파일 저장소가 삭제되었습니다!")
                                if 'delete_storage_id' in st.session_state:
                                    del st.session_state.delete_storage_id
                                st.rerun()
                            else:
                                st.error("파일 저장소 삭제에 실패했습니다.")
                    with col_del_2:
                        if st.button("❌ 취소", key=f"cancel_delete_{storage['file_id']}"):
                            if 'delete_storage_id' in st.session_state:
                                del st.session_state.delete_storage_id
                            st.rerun()
    
    elif menu == "🔍 고급 검색":
        st.header("🔍 파일 저장소 고급 검색")
        
        # 세션 상태 초기화
        if 'search_results' not in st.session_state:
            st.session_state.search_results = None
        if 'search_performed' not in st.session_state:
            st.session_state.search_performed = False
        
        # 검색 옵션
        col1, col2, col3 = st.columns([2, 2, 1])
        
        with col1:
            search_term = st.text_input(
                "🔎 검색어", 
                placeholder="제목, 작성자, 개요에서 검색...",
                value=st.session_state.get('last_search_term', ''),
                key="search_term_input"
            )
            all_authors = get_all_authors()
            author_options = ["전체"] + all_authors
            
            # 저장된 작성자 필터가 목록에 있는지 확인
            saved_author = st.session_state.get('last_author_filter', '전체')
            try:
                author_index = author_options.index(saved_author)
            except ValueError:
                author_index = 0  # 목록에 없으면 "전체"로 설정
            
            author_filter = st.selectbox(
                "👤 작성자 필터", 
                author_options,
                index=author_index
            )
        
        with col2:
            date_from = st.date_input(
                "📅 시작일", 
                value=st.session_state.get('last_date_from', None),
                key="date_from_input"
            )
            date_to = st.date_input(
                "📅 종료일", 
                value=st.session_state.get('last_date_to', None),
                key="date_to_input"
            )
        
        with col3:
            st.write("")  # 빈 공간
            search_clicked = st.button("🔍 검색", type="primary", use_container_width=True)
            reset_clicked = st.button("🔄 초기화", use_container_width=True)
        
        # 초기화 처리
        if reset_clicked:
            st.session_state.search_results = None
            st.session_state.search_performed = False
            if 'last_search_term' in st.session_state:
                del st.session_state.last_search_term
            if 'last_author_filter' in st.session_state:
                del st.session_state.last_author_filter
            if 'last_date_from' in st.session_state:
                del st.session_state.last_date_from
            if 'last_date_to' in st.session_state:
                del st.session_state.last_date_to
            st.rerun()
        
        # 검색 실행
        if search_clicked:
            # 검색 조건을 세션에 저장
            st.session_state.last_search_term = search_term
            st.session_state.last_author_filter = author_filter
            st.session_state.last_date_from = date_from
            st.session_state.last_date_to = date_to
            
            # 검색 실행
            storages = search_file_storages(
                search_term=search_term if search_term else None,
                author_filter=author_filter,
                date_from=date_from,
                date_to=date_to
            )
            
            # 검색 결과를 세션에 저장
            st.session_state.search_results = storages
            st.session_state.search_performed = True
        
        # 전체 화면 파일 미리보기 (고급 검색용)
        if st.session_state.get('preview_file_id'):
            st.markdown("---")
            st.markdown("## 🔍 파일 미리보기 (전체 화면)")
            
            col_preview_1, col_preview_2 = st.columns([10, 1])
            with col_preview_2:
                if st.button("❌ 닫기", key="close_search_preview"):
                    del st.session_state.preview_file_id
                    del st.session_state.preview_filename
                    del st.session_state.preview_file_type
                    st.rerun()
            
            with col_preview_1:
                file_data = get_file_binary_data(st.session_state.preview_file_id)
                if file_data:
                    display_file_preview(
                        file_data, 
                        st.session_state.preview_file_type, 
                        st.session_state.preview_filename
                    )
                else:
                    st.error("파일 데이터를 불러올 수 없습니다.")
            
            st.markdown("---")

        # 검색 결과 표시
        if st.session_state.search_performed and st.session_state.search_results is not None:
            storages = st.session_state.search_results
            
            # 검색 조건 요약 표시
            search_summary = []
            if st.session_state.get('last_search_term'):
                search_summary.append(f"검색어: '{st.session_state.last_search_term}'")
            if st.session_state.get('last_author_filter') and st.session_state.last_author_filter != "전체":
                search_summary.append(f"작성자: {st.session_state.last_author_filter}")
            if st.session_state.get('last_date_from'):
                search_summary.append(f"시작일: {st.session_state.last_date_from}")
            if st.session_state.get('last_date_to'):
                search_summary.append(f"종료일: {st.session_state.last_date_to}")
            
            if search_summary:
                st.info(f"🔍 **검색 조건**: {' | '.join(search_summary)}")
            
            if storages:
                st.success(f"**검색 결과: {len(storages)}개**")
                
                # 선택된 저장소 상세 보기
                if st.session_state.get('selected_storage'):
                    selected_storage = None
                    for storage in storages:
                        if storage['file_id'] == st.session_state.selected_storage:
                            selected_storage = storage
                            break
                    
                    if selected_storage:
                        st.markdown("---")
                        st.markdown(f"## 📂 {selected_storage['title']} 상세 보기")
                        
                        col_back, col_title = st.columns([1, 10])
                        with col_back:
                            if st.button("← 뒤로", key="back_to_search"):
                                del st.session_state.selected_storage
                                st.rerun()
                        
                        # 파일 목록 표시
                        files = get_storage_files(selected_storage['file_id'])
                        if files:
                            st.subheader("📄 파일 목록")
                            for file_info in files:
                                col1, col2, col3 = st.columns([2, 1, 1])
                                
                                with col1:
                                    st.write(f"**{file_info['filename']}** ({file_info['file_type'].upper()})")
                                    st.caption(f"크기: {file_info['file_size']:,} bytes")
                                
                                with col2:
                                    if st.button(f"👁️ 미리보기", key=f"search_preview_{file_info['file_data_id']}"):
                                        st.session_state.preview_file_id = file_info['file_data_id']
                                        st.session_state.preview_filename = file_info['filename']
                                        st.session_state.preview_file_type = file_info['file_type']
                                
                                with col3:
                                    file_data = get_file_binary_data(file_info['file_data_id'])
                                    if file_data:
                                        mime_type = get_file_mime_type(file_data['file_type'])
                                        st.download_button(
                                            label="💾 다운로드",
                                            data=file_data['binary_data'],
                                            file_name=file_data['filename'],
                                            mime=mime_type,
                                            key=f"search_download_{file_info['file_data_id']}"
                                        )
                        
                        # AI 분석 결과 표시
                        analyses = get_ai_analysis(selected_storage['file_id'])
                        if analyses:
                            st.subheader("🤖 AI 분석 결과")
                            for analysis in analyses:
                                st.write(f"**모델:** {analysis['model_name']}")
                                st.write(f"**분석 시간:** {analysis['created_at'].strftime('%Y-%m-%d %H:%M')}")
                                st.markdown(analysis['analysis_content'])
                                st.divider()
                        
                        st.markdown("---")
                
                for storage in storages:
                    with st.container():
                        col1, col2 = st.columns([4, 1])
                        
                        with col1:
                            st.write(f"**📁 {storage['title']}**")
                            st.write(f"작성자: {storage['author']} | 파일 수: {storage['actual_file_count']} | 분석: {storage['analysis_count']}개")
                            st.write(f"생성일: {storage['created_at'].strftime('%Y-%m-%d %H:%M')}")
                            if storage['description']:
                                st.write(f"개요: {storage['description']}")
                        
                        with col2:
                            col2_1, col2_2, col2_3 = st.columns(3)
                            with col2_1:
                                if st.button("📂", key=f"search_view_{storage['file_id']}", help="상세 보기"):
                                    st.session_state.selected_storage = storage['file_id']
                                    st.rerun()
                            with col2_2:
                                if st.button("✏️", key=f"search_edit_{storage['file_id']}", help="편집"):
                                    st.session_state.edit_storage_id = storage['file_id']
                                    st.session_state.return_to_search = True  # 검색으로 돌아가기 플래그
                                    st.rerun()
                            with col2_3:
                                if st.button("🗑️", key=f"search_delete_{storage['file_id']}", help="삭제"):
                                    st.session_state.delete_storage_id = storage['file_id']
                                    st.rerun()
                        
                        # 삭제 확인 대화상자 (검색 결과에서도 표시)
                        if st.session_state.get('delete_storage_id') == storage['file_id']:
                            st.error("⚠️ 정말로 이 파일 저장소를 삭제하시겠습니까?")
                            st.write("**삭제될 내용:**")
                            st.write(f"- 저장소: {storage['title']}")
                            st.write(f"- 파일: {storage['actual_file_count']}개")
                            
                            # 분석 결과 개수 확인
                            if storage['analysis_count'] > 0:
                                st.write(f"- AI 분석 결과: {storage['analysis_count']}개")
                            
                            col_del_1, col_del_2 = st.columns(2)
                            with col_del_1:
                                if st.button("✅ 삭제 확인", key=f"search_confirm_delete_{storage['file_id']}", type="primary"):
                                    if delete_file_storage(storage['file_id']):
                                        st.success("파일 저장소가 삭제되었습니다!")
                                        # 검색 결과 새로고침
                                        if 'delete_storage_id' in st.session_state:
                                            del st.session_state.delete_storage_id
                                        # 검색 다시 실행
                                        storages = search_file_storages(
                                            search_term=st.session_state.get('last_search_term'),
                                            author_filter=st.session_state.get('last_author_filter'),
                                            date_from=st.session_state.get('last_date_from'),
                                            date_to=st.session_state.get('last_date_to')
                                        )
                                        st.session_state.search_results = storages
                                        st.rerun()
                                    else:
                                        st.error("파일 저장소 삭제에 실패했습니다.")
                            with col_del_2:
                                if st.button("❌ 취소", key=f"search_cancel_delete_{storage['file_id']}"):
                                    if 'delete_storage_id' in st.session_state:
                                        del st.session_state.delete_storage_id
                                    st.rerun()
                        
                        st.divider()
            else:
                st.warning("검색 결과가 없습니다.")
        elif st.session_state.search_performed:
            st.info("검색을 실행해주세요.")
    
    elif menu == "✏️ 저장소 편집":
        st.header("✏️ 파일 저장소 편집")
        
        # 편집할 저장소 선택
        if 'edit_storage_id' not in st.session_state:
            storages = get_file_storages()
            
            if not storages:
                st.info("편집할 파일 저장소가 없습니다.")
                return
            
            st.subheader("편집할 저장소 선택")
            for storage in storages:
                col1, col2 = st.columns([4, 1])
                
                with col1:
                    st.write(f"**📁 {storage['title']}** ({storage['author']})")
                    st.caption(f"파일 {storage['file_count']}개 | {storage['created_at'].strftime('%Y-%m-%d %H:%M')}")
                
                with col2:
                    if st.button("✏️ 편집", key=f"edit_select_{storage['file_id']}"):
                        st.session_state.edit_storage_id = storage['file_id']
                        st.rerun()
        else:
            # 선택된 저장소 편집
            storage_id = st.session_state.edit_storage_id
            storage_data = get_file_storage_by_id(storage_id)
            
            if not storage_data:
                st.error("저장소를 찾을 수 없습니다.")
                if st.button("목록으로 돌아가기"):
                    del st.session_state.edit_storage_id
                    st.rerun()
                return
            
            st.subheader(f"📁 {storage_data['title']} 편집")
            
            with st.form("edit_storage_form"):
                col1, col2 = st.columns(2)
                
                with col1:
                    title = st.text_input("📝 제목", value=storage_data['title'])
                    author = st.text_input("✍️ 작성자", value=storage_data['author'])
                
                with col2:
                    description = st.text_area("📋 개요", value=storage_data['description'], height=100)
                
                col_btn1, col_btn2 = st.columns(2)
                with col_btn1:
                    submitted = st.form_submit_button("💾 수정 저장", type="primary")
                with col_btn2:
                    cancelled = st.form_submit_button("❌ 취소")
                
                if submitted:
                    if not title:
                        st.error("제목을 입력해 주세요.")
                    else:
                        update_data = {
                            'title': title,
                            'author': author or "Unknown",
                            'description': description
                        }
                        
                        if update_file_storage(storage_id, update_data):
                            st.success("✅ 파일 저장소가 성공적으로 수정되었습니다!")
                            
                            # 검색 결과 새로고침 (검색에서 온 경우)
                            if st.session_state.get('return_to_search') and st.session_state.get('search_performed'):
                                storages = search_file_storages(
                                    search_term=st.session_state.get('last_search_term'),
                                    author_filter=st.session_state.get('last_author_filter'),
                                    date_from=st.session_state.get('last_date_from'),
                                    date_to=st.session_state.get('last_date_to')
                                )
                                st.session_state.search_results = storages
                            
                            # 세션 정리
                            del st.session_state.edit_storage_id
                            if 'return_to_search' in st.session_state:
                                del st.session_state.return_to_search
                            
                            time.sleep(1)
                            st.rerun()
                        else:
                            st.error("❌ 수정에 실패했습니다.")
                
                if cancelled:
                    del st.session_state.edit_storage_id
                    if 'return_to_search' in st.session_state:
                        del st.session_state.return_to_search
                    st.rerun()
            
            # 현재 저장소 정보 표시
            st.subheader("📊 현재 저장소 정보")
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.metric("파일 수", storage_data['file_count'])
            with col2:
                st.metric("생성일", storage_data['created_at'].strftime('%Y-%m-%d'))
            with col3:
                # AI 분석 개수 조회
                analyses = get_ai_analysis(storage_id)
                st.metric("AI 분석", len(analyses))
            
            # 관련 파일 목록 표시 및 관리
            files = get_storage_files(storage_id)
            if files:
                st.subheader("📄 파일 관리")
                for file_info in files:
                    col_file_1, col_file_2 = st.columns([4, 1])
                    
                    with col_file_1:
                        st.write(f"• **{file_info['filename']}** ({file_info['file_type'].upper()}) - {file_info['file_size']:,} bytes")
                    
                    with col_file_2:
                        if st.button("🔄 교체", key=f"edit_replace_btn_{file_info['file_data_id']}", use_container_width=True):
                            st.session_state[f'edit_replace_file_{file_info["file_data_id"]}'] = True
                        
                        if st.button("🗑️ 삭제", key=f"edit_delete_{file_info['file_data_id']}", use_container_width=True):
                            if delete_storage_file(file_info['file_data_id']):
                                st.success("파일이 삭제되었습니다!")
                                st.rerun()
                            else:
                                st.error("파일 삭제에 실패했습니다.")
                    
                    # 파일 교체 UI
                    if st.session_state.get(f'edit_replace_file_{file_info["file_data_id"]}'):
                        st.write("**파일 교체**")
                        replacement_file = st.file_uploader(
                            f"새 파일 선택 (현재: {file_info['filename']})",
                            type=['pdf', 'docx', 'pptx', 'xlsx', 'md', 'txt', 'jpg', 'jpeg', 'png'],
                            key=f"edit_replace_upload_{file_info['file_data_id']}"
                        )
                        
                        col_rep_1, col_rep_2 = st.columns(2)
                        with col_rep_1:
                            if st.button("✅ 교체 확인", key=f"edit_confirm_replace_{file_info['file_data_id']}"):
                                if replacement_file:
                                    new_file_data = parse_uploaded_file(replacement_file)
                                    if new_file_data and replace_storage_file(file_info['file_data_id'], new_file_data):
                                        st.success("파일이 교체되었습니다!")
                                        del st.session_state[f'edit_replace_file_{file_info["file_data_id"]}']
                                        st.rerun()
                                    else:
                                        st.error("파일 교체에 실패했습니다.")
                                else:
                                    st.error("교체할 파일을 선택해 주세요.")
                        with col_rep_2:
                            if st.button("❌ 취소", key=f"edit_cancel_replace_{file_info['file_data_id']}"):
                                del st.session_state[f'edit_replace_file_{file_info["file_data_id"]}']
                                st.rerun()
                
                # 새 파일 추가
                st.subheader("➕ 새 파일 추가")
                new_files = st.file_uploader(
                    "새 파일 선택",
                    type=['pdf', 'docx', 'pptx', 'xlsx', 'md', 'txt', 'jpg', 'jpeg', 'png'],
                    accept_multiple_files=True,
                    key=f"edit_add_files_{storage_id}"
                )
                
                if new_files and st.button("➕ 파일 추가", key=f"edit_add_files_btn_{storage_id}"):
                    with st.spinner("파일을 추가하는 중..."):
                        new_files_data = []
                        for new_file in new_files:
                            file_data = parse_uploaded_file(new_file)
                            if file_data:
                                new_files_data.append(file_data)
                        
                        if new_files_data and add_files_to_storage(storage_id, new_files_data):
                            st.success(f"{len(new_files_data)}개 파일이 추가되었습니다!")
                            st.rerun()
                        else:
                            st.error("파일 추가에 실패했습니다.")
            else:
                st.info("포함된 파일이 없습니다.")
                
                # 파일이 없는 경우에도 새 파일 추가 기능 제공
                st.subheader("➕ 파일 추가")
                new_files = st.file_uploader(
                    "파일 선택",
                    type=['pdf', 'docx', 'pptx', 'xlsx', 'md', 'txt', 'jpg', 'jpeg', 'png'],
                    accept_multiple_files=True,
                    key=f"edit_add_files_empty_{storage_id}"
                )
                
                if new_files and st.button("➕ 파일 추가", key=f"edit_add_files_empty_btn_{storage_id}"):
                    with st.spinner("파일을 추가하는 중..."):
                        new_files_data = []
                        for new_file in new_files:
                            file_data = parse_uploaded_file(new_file)
                            if file_data:
                                new_files_data.append(file_data)
                        
                        if new_files_data and add_files_to_storage(storage_id, new_files_data):
                            st.success(f"{len(new_files_data)}개 파일이 추가되었습니다!")
                            st.rerun()
                        else:
                            st.error("파일 추가에 실패했습니다.")
    
    elif menu == "📊 통계":
        st.header("📊 파일 저장소 통계")
        
        storages = get_file_storages()
        
        if storages:
            df = pd.DataFrame(storages)
            
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("전체 저장소", len(storages))
            with col2:
                total_files = df['file_count'].sum()
                st.metric("전체 파일", total_files)
            with col3:
                unique_authors = df['author'].nunique()
                st.metric("작성자 수", unique_authors)
            
            # 시간별 생성 통계
            df['date'] = pd.to_datetime(df['created_at']).dt.date
            date_counts = df.groupby('date').size().reset_index(name='count')
            
            fig = px.line(date_counts, x='date', y='count', 
                         title='📈 일별 저장소 생성 현황')
            st.plotly_chart(fig, use_container_width=True)
            
            # 작성자별 통계
            author_counts = df['author'].value_counts().head(10)
            fig2 = px.bar(x=author_counts.index, y=author_counts.values,
                         title='👥 작성자별 저장소 수')
            st.plotly_chart(fig2, use_container_width=True)
        else:
            st.info("통계를 표시할 데이터가 없습니다.")

if __name__ == "__main__":
    main() 